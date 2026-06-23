import copy
from datetime import datetime
import os
import re
import sys
import warnings

warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl.reader.drawings")

try:
    import torch
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 torch，请先在 Colab 执行: pip install torch") from exc

try:
    import pandas as pd
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 pandas，请先在 Colab 执行: pip install pandas openpyxl xlrd") from exc

try:
    from tqdm.auto import tqdm
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 tqdm，请先在 Colab 执行: pip install tqdm") from exc

try:
    from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 transformers，请先在 Colab 执行: pip install transformers sentencepiece") from exc

try:
    from peft import PeftModel
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 peft，请先在 Colab 执行: pip install peft accelerate") from exc

def _version_lt(v: str, target: str) -> bool:
    def _to_tuple(s: str):
        parts = []
        for x in str(s).split("."):
            num = ""
            for ch in x:
                if ch.isdigit():
                    num += ch
                else:
                    break
            parts.append(int(num) if num else 0)
        return tuple(parts)

    return _to_tuple(v) < _to_tuple(target)


def disable_incompatible_torchao_for_peft():
    try:
        import importlib.metadata as _md

        torchao_ver = _md.version("torchao")
    except Exception:
        torchao_ver = ""

    try:
        import peft.import_utils as _peft_import_utils
    except Exception:
        return

    if not hasattr(_peft_import_utils, "_orig_is_torchao_available"):
        _peft_import_utils._orig_is_torchao_available = _peft_import_utils.is_torchao_available

        def _safe_is_torchao_available():
            try:
                return _peft_import_utils._orig_is_torchao_available()
            except ImportError:
                return False

        _peft_import_utils.is_torchao_available = _safe_is_torchao_available

    if torchao_ver and _version_lt(torchao_ver, "0.16.0"):
        print(
            f"[兼容修复] 检测到不兼容的 torchao {torchao_ver}，已自动禁用 torchao 分支，避免 PEFT LoRA 注入失败。",
            flush=True,
        )


disable_incompatible_torchao_for_peft()

try:
    from huggingface_hub import snapshot_download
    from huggingface_hub.utils import enable_progress_bars
except Exception:
    snapshot_download = None
    enable_progress_bars = None

try:
    from pptx import Presentation
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 python-pptx，请先在 Colab 执行: pip install python-pptx") from exc

try:
    from openpyxl import load_workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 openpyxl，请先在 Colab 执行: pip install openpyxl") from exc

try:
    from docx import Document
    from docx.enum.text import WD_BREAK
    from docx.text.paragraph import Paragraph
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
except ModuleNotFoundError as exc:
    raise SystemExit("缺少依赖 python-docx，请先在 Colab 执行: pip install python-docx") from exc


BASE_MODEL_NAME = (os.getenv("NLLB_BASE_MODEL") or "facebook/nllb-200-3.3B").strip()
ADAPTER_DIR = (os.getenv("NLLB_ADAPTER_DIR") or "/content/drive/MyDrive/translated_models/NLLB200-33B/adapter_checkpoint4000").strip()
MERGE_ADAPTER = int(os.getenv("NLLB_MERGE_ADAPTER", "0") or "0")
HF_CACHE_DIR = (os.getenv("HF_CACHE_DIR") or "/content/hf_cache").strip()

DEFAULT_INPUT_FILE = (os.getenv("DEFAULT_INPUT_FILE") or "/content/drive/MyDrive/File_translate/corpus.xlsx").strip()
DEFAULT_SRC_LANG = (os.getenv("NLLB_SRC_LANG") or "kor_Hang").strip()
DEFAULT_TGT_LANG = (os.getenv("NLLB_TGT_LANG") or "zho_Hans").strip()

NLLB_DEFAULT_BATCH_SIZE = int(os.getenv("NLLB_BATCH_SIZE", "32") or "32")
NLLB_DEFAULT_MAX_LENGTH = int(os.getenv("NLLB_MAX_LENGTH", "96") or "96")
NLLB_MAX_SOURCE_LENGTH = int(os.getenv("NLLB_MAX_SOURCE_LENGTH", "256") or "256")
NLLB_NUM_BEAMS = int(os.getenv("NLLB_NUM_BEAMS", "1") or "1")
NLLB_USE_CACHE = int(os.getenv("NLLB_USE_CACHE", "1") or "1")
APPEND_TRANSLATION = int(os.getenv("APPEND_TRANSLATION", "0") or "0")
GENERATE_CORPUS = int(os.getenv("GENERATE_CORPUS", "0") or "0")
CORPUS_DIR = (os.getenv("CORPUS_DIR") or "/content/drive/MyDrive/File_translate/Corpus").strip()
TARGET_FONT_NAME = (os.getenv("NLLB_TARGET_FONT") or "Microsoft YaHei").strip()
USE_TQDM = int(os.getenv("USE_TQDM", "0") or "0")
PROGRESS_PRINT_EVERY = int(os.getenv("PROGRESS_PRINT_EVERY", "20") or "20")
PROTECT_TOKENS = int(os.getenv("PROTECT_TOKENS", "1") or "1")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
HANGUL_RE = re.compile(r"[\uac00-\ud7af\u1100-\u11ff\u3130-\u318f\ua960-\ua97f\ud7b0-\ud7ff]")
CHINESE_RE = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf\uf900-\ufaff]")

_translation_cache: dict[tuple[str, str, str], str] = {}
_nllb_tokenizer = None
_nllb_model = None
_nllb_device = "cuda" if torch.cuda.is_available() else "cpu"
_nllb_loaded_signature = None
original_texts: list[str] = []
translated_texts: list[str] = []


def reset_runtime_buffers():
    original_texts.clear()
    translated_texts.clear()


def save_to_corpus(orig: list[str], trans: list[str]):
    if not GENERATE_CORPUS or not orig:
        return None

    os.makedirs(CORPUS_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    direction = f"{DEFAULT_SRC_LANG}2{DEFAULT_TGT_LANG}"
    corpus_file = os.path.join(CORPUS_DIR, f"Corpus_nllb33b_{direction}_{timestamp}.xlsx")

    seen = set()
    filtered_orig = []
    filtered_trans = []
    target_is_chinese = str(DEFAULT_TGT_LANG).startswith("zho")

    for o, t in zip(orig, trans):
        if o is None or t is None:
            continue
        o = _normalize_text_newlines(str(o)).strip()
        t = _normalize_text_newlines(str(t)).strip()
        if not o or not t:
            continue
        if o == t:
            continue

        has_korean = bool(HANGUL_RE.search(o))
        has_chinese = bool(CHINESE_RE.search(o))

        if has_korean:
            pass
        else:
            if has_chinese and target_is_chinese:
                continue
            if not has_chinese:
                continue

        key = o
        if key in seen:
            continue
        seen.add(key)
        filtered_orig.append(o)
        filtered_trans.append(t)

    if not filtered_orig:
        return None

    pd.DataFrame(
        {
            "序号": range(1, len(filtered_orig) + 1),
            "翻译前": filtered_orig,
            "翻译后": filtered_trans,
        }
    ).to_excel(corpus_file, index=False)
    print(f"[完成] 语料库已保存: {corpus_file}")
    return corpus_file


def _adapter_present(adapter_dir: str) -> bool:
    if not adapter_dir or not os.path.isdir(adapter_dir):
        return False
    try:
        names = set(os.listdir(adapter_dir))
    except Exception:
        return False
    return (
        "adapter_config.json" in names
        and ("adapter_model.safetensors" in names or "adapter_model.bin" in names)
    )


def _maybe_snapshot_download(model_id_or_dir: str) -> str:
    model_id_or_dir = str(model_id_or_dir or "").strip()
    if not model_id_or_dir:
        raise RuntimeError("base_model_name 不能为空")
    if os.path.isdir(model_id_or_dir):
        return model_id_or_dir
    if snapshot_download is None:
        return model_id_or_dir
    if enable_progress_bars is not None:
        try:
            enable_progress_bars()
        except Exception:
            pass
    os.makedirs(HF_CACHE_DIR, exist_ok=True)
    try:
        local_dir = snapshot_download(
            repo_id=model_id_or_dir,
            cache_dir=HF_CACHE_DIR,
            local_files_only=True,
        )
        print(f"[系统] 已命中本地 HuggingFace 缓存，无需重新下载: {local_dir}")
        return local_dir
    except Exception:
        pass

    print("[系统] 本地缓存不完整，开始从 HuggingFace 下载/补齐缺失文件（无账号也可下载公开模型，速度取决于网络）...")
    local_dir = snapshot_download(
        repo_id=model_id_or_dir,
        cache_dir=HF_CACHE_DIR,
        resume_download=True,
    )
    print(f"[系统] HuggingFace 缓存目录: {local_dir}")
    return local_dir


def load_nllb_model(base_model_name: str | None = None, adapter_dir: str | None = None):
    global _nllb_tokenizer, _nllb_model, _nllb_loaded_signature

    base_model_name = (base_model_name or BASE_MODEL_NAME).strip()
    adapter_dir = (adapter_dir or ADAPTER_DIR).strip()
    signature = (base_model_name, adapter_dir, int(MERGE_ADAPTER), _nllb_device)
    if _nllb_model is not None and _nllb_loaded_signature == signature:
        return _nllb_tokenizer, _nllb_model

    print(f"[系统] Base model: {base_model_name}")
    print(f"[系统] Adapter dir: {adapter_dir}")
    print(f"[系统] 使用设备: {_nllb_device}")
    if _nllb_device == "cuda":
        print(f"[系统] GPU: {torch.cuda.get_device_name(0)}")
        print(f"[系统] GPU 显存: {torch.cuda.get_device_properties(0).total_memory / 1024**3:.1f} GB")

    base_model_resolved = _maybe_snapshot_download(base_model_name)
    tokenizer = AutoTokenizer.from_pretrained(base_model_resolved, local_files_only=os.path.isdir(base_model_resolved))
    torch_dtype = torch.float16 if _nllb_device == "cuda" else torch.float32
    try:
        model = AutoModelForSeq2SeqLM.from_pretrained(
            base_model_resolved,
            dtype=torch_dtype,
            low_cpu_mem_usage=True,
            local_files_only=os.path.isdir(base_model_resolved),
        )
    except TypeError:
        model = AutoModelForSeq2SeqLM.from_pretrained(
            base_model_resolved,
            torch_dtype=torch_dtype,
            low_cpu_mem_usage=True,
            local_files_only=os.path.isdir(base_model_resolved),
        )
    model = model.to(_nllb_device)
    model.eval()
    try:
        model.config.use_cache = bool(NLLB_USE_CACHE)
    except Exception:
        pass

    if _adapter_present(adapter_dir):
        print("[系统] 检测到 adapter，加载并应用到 base model。")
        model = PeftModel.from_pretrained(model, adapter_dir)
        if MERGE_ADAPTER:
            print("[系统] 正在合并 adapter 到 base model（merge_and_unload）...")
            model = model.merge_and_unload()
        model.eval()
    else:
        print("[系统] 未检测到 adapter（或目录为空），将使用 base model 直接翻译。")

    _nllb_tokenizer = tokenizer
    _nllb_model = model
    _nllb_loaded_signature = signature
    print("[系统] 模型加载完成。")
    return tokenizer, model


def _normalize_text_newlines(text: str) -> str:
    return str(text or "").replace("\r\n", "\n").replace("\r", "\n")


def _ppt_normalize_linebreaks(text: str) -> str:
    return _normalize_text_newlines(str(text or "")).replace("\v", "\n")


def _ppt_denormalize_linebreaks(text: str) -> str:
    return str(text or "").replace("\n", "\v")


def _text_contains_korean(text: str) -> bool:
    return bool(HANGUL_RE.search(str(text or "")))


def _text_contains_english(text: str) -> bool:
    return bool(re.search(r"[A-Za-z]", str(text or "")))



def _looks_like_excel_formula(text: str) -> bool:
    t = str(text or "").lstrip()
    return bool(t) and t.startswith("=")

_PROTECT_RE = re.compile(
    r"("
    r"https?://\S+"
    r"|[\w\.-]+@[\w\.-]+\.\w+"
    r"|[$¥₩€£]\s*(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?"
    r"|(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?(?:\s*[A-Za-z%°℃]+)?"
    r"|\b[A-Za-z0-9][A-Za-z0-9._/-]*\d[A-Za-z0-9._/-]*\b"
    r")"
)


def _protect_tokens(text: str):
    if not PROTECT_TOKENS:
        return str(text or ""), {}
    mapping = {}
    idx = 0

    def _repl(m):
        nonlocal idx
        key = f"__K{idx}__"
        idx += 1
        mapping[key] = m.group(0)
        return key

    protected = _PROTECT_RE.sub(_repl, str(text or ""))
    return protected, mapping


def _restore_tokens(text: str, mapping: dict[str, str]) -> str:
    out = str(text or "")
    for k, v in mapping.items():
        out = out.replace(k, v)
    return out


def _translate_text_with_protection(text: str) -> str:
    protected, mapping = _protect_tokens(text)
    translated = get_translation(protected)
    return _restore_tokens(translated, mapping)


def _translate_texts_with_protection(texts: list[str]) -> list[str]:
    protected_list = []
    mappings = []
    for t in texts:
        protected, mapping = _protect_tokens(t)
        protected_list.append(protected)
        mappings.append(mapping)
    translated_list = get_translation_batch(protected_list)
    out = []
    for translated, mapping in zip(translated_list, mappings):
        out.append(_restore_tokens(translated, mapping))
    return out


def _translate_preserve_newlines(text: str) -> str:
    normalized = _normalize_text_newlines(str(text or ""))
    if "\n" not in normalized:
        return get_translation(normalized)
    lines = normalized.split("\n")
    line_pos = []
    line_texts = []
    for i, line in enumerate(lines):
        if line.strip() and _text_contains_korean(line):
            line_pos.append(i)
            line_texts.append(line)
    if not line_texts:
        return get_translation(normalized)
    translated = get_translation_batch(line_texts)
    for i, t in zip(line_pos, translated):
        lines[i] = t
    return "\n".join(lines)


def _translate_preserve_newlines_with_protection(text: str, allow_english: bool) -> str:
    normalized = _normalize_text_newlines(str(text or ""))
    if "\n" not in normalized:
        if allow_english and (_text_contains_korean(normalized) or _text_contains_english(normalized)):
            return _translate_text_with_protection(normalized)
        if _text_contains_korean(normalized):
            return _translate_text_with_protection(normalized)
        return normalized

    lines = normalized.split("\n")
    line_pos = []
    line_texts = []
    for i, line in enumerate(lines):
        if not line.strip():
            continue
        if _text_contains_korean(line) or (allow_english and _text_contains_english(line)):
            line_pos.append(i)
            line_texts.append(line)
    if not line_texts:
        return normalized
    translated = _translate_texts_with_protection(line_texts)
    for i, t in zip(line_pos, translated):
        lines[i] = t
    return "\n".join(lines)


def _translate_batch_core(text_list: list[str], src_lang: str = DEFAULT_SRC_LANG, tgt_lang: str = DEFAULT_TGT_LANG) -> list[str]:
    if not text_list:
        return []
    tokenizer, model = load_nllb_model()
    batch_size = int(NLLB_DEFAULT_BATCH_SIZE)
    max_length = int(NLLB_DEFAULT_MAX_LENGTH)
    max_source_length = int(NLLB_MAX_SOURCE_LENGTH)

    lang_code_to_id = getattr(tokenizer, "lang_code_to_id", None)
    if isinstance(lang_code_to_id, dict) and tgt_lang in lang_code_to_id:
        forced_bos_token_id = int(lang_code_to_id[tgt_lang])
    else:
        forced_bos_token_id = tokenizer.convert_tokens_to_ids(tgt_lang)

    outs = []
    total_batches = (len(text_list) + batch_size - 1) // batch_size
    starts = range(0, len(text_list), batch_size)
    can_tqdm = bool(USE_TQDM) and (getattr(sys.stdout, "isatty", lambda: False)() or getattr(sys.stderr, "isatty", lambda: False)())
    if can_tqdm:
        starts_iter = tqdm(starts, desc="翻译进度", total=total_batches)
    else:
        starts_iter = starts
        if total_batches > 1:
            print(f"[进度] 批量翻译: 0/{total_batches}", flush=True)

    done_batches = 0
    for start in starts_iter:
        batch = text_list[start : start + batch_size]
        tokenizer.src_lang = src_lang
        model_max_len = int(getattr(tokenizer, "model_max_length", 512) or 512)
        src_max_len = min(int(max_source_length), int(model_max_len))
        encoded = tokenizer(
            batch,
            return_tensors="pt",
            padding=True,
            truncation=True,
            max_length=src_max_len,
        )
        encoded = {k: v.to(_nllb_device) for k, v in encoded.items()}
        with torch.inference_mode():
            generated = model.generate(
                **encoded,
                forced_bos_token_id=int(forced_bos_token_id),
                max_length=max_length,
                use_cache=bool(NLLB_USE_CACHE),
                do_sample=False,
                num_beams=int(NLLB_NUM_BEAMS),
                temperature=1.0,
                top_p=1.0,
                pad_token_id=tokenizer.pad_token_id,
                eos_token_id=tokenizer.eos_token_id,
            )
        decoded = tokenizer.batch_decode(generated, skip_special_tokens=True)
        outs.extend([("" if x is None else str(x).strip()) for x in decoded])
        if not can_tqdm and total_batches > 1:
            done_batches += 1
            if done_batches == total_batches or (PROGRESS_PRINT_EVERY > 0 and (done_batches % PROGRESS_PRINT_EVERY == 0)):
                print(f"[进度] 批量翻译: {done_batches}/{total_batches}", flush=True)
    return outs


def get_translation(text: str, src_lang: str = DEFAULT_SRC_LANG, tgt_lang: str = DEFAULT_TGT_LANG) -> str:
    text = _normalize_text_newlines(text).strip()
    if not text:
        return text
    cache_key = (src_lang, tgt_lang, text)
    cached = _translation_cache.get(cache_key)
    if cached is not None:
        return cached
    translated = _translate_batch_core([text], src_lang=src_lang, tgt_lang=tgt_lang)[0]
    _translation_cache[cache_key] = translated
    return translated


def get_translation_batch(text_list: list[str], src_lang: str = DEFAULT_SRC_LANG, tgt_lang: str = DEFAULT_TGT_LANG) -> list[str]:
    if not text_list:
        return []
    normalized = [_normalize_text_newlines(x).strip() for x in text_list]
    outputs = [""] * len(normalized)
    missing_indices = []
    missing_texts = []
    for idx, text in enumerate(normalized):
        if not text:
            outputs[idx] = text
            continue
        cache_key = (src_lang, tgt_lang, text)
        cached = _translation_cache.get(cache_key)
        if cached is not None:
            outputs[idx] = cached
        else:
            missing_indices.append(idx)
            missing_texts.append(text)
    if missing_texts:
        translated_missing = _translate_batch_core(missing_texts, src_lang=src_lang, tgt_lang=tgt_lang)
        for idx, src_text, trans in zip(missing_indices, missing_texts, translated_missing):
            outputs[idx] = trans
            _translation_cache[(src_lang, tgt_lang, src_text)] = trans
    return outputs


def append_translation_to_original(text: str, translated_text: str, cell=None) -> str:
    text = str(text or "").strip()
    translated_text = str(translated_text or "").strip()
    result = f"{text}\n{translated_text}" if text and translated_text else (text or translated_text)
    if cell is not None:
        cell.alignment = Alignment(wrap_text=True, vertical="center")
        ws = cell.parent
        row_num = cell.row
        original_height = ws.row_dimensions[row_num].height
        ws.row_dimensions[row_num].height = (original_height * 2) if original_height else 30
    return result


def clean_sheet_name(name: str) -> str:
    if not name:
        return "Sheet"
    invalid_chars = r'[\\/?:*\[\](){}<>|"\']'
    return re.sub(invalid_chars, "", str(name))[:31]


def get_target_font_name() -> str:
    return TARGET_FONT_NAME


def set_docx_r_element_font(r_element, font_name: str):
    if r_element is None or not font_name:
        return
    rPr = r_element.get_or_add_rPr()
    rFonts = rPr.rFonts
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), font_name)
    rFonts.set(qn("w:hAnsi"), font_name)
    rFonts.set(qn("w:eastAsia"), font_name)
    rFonts.set(qn("w:cs"), font_name)
    for k in (qn("w:asciiTheme"), qn("w:hAnsiTheme"), qn("w:eastAsiaTheme"), qn("w:csTheme")):
        try:
            rFonts.attrib.pop(k, None)
        except Exception:
            pass


def set_docx_run_font(run, font_name: str):
    if not run or not font_name:
        return
    run.font.name = font_name
    set_docx_r_element_font(run._r, font_name)


def apply_run_format(source_run, target_run):
    if not source_run or not target_run:
        return
    source_rPr = source_run._element.rPr
    if source_rPr is not None:
        new_rPr = copy.deepcopy(source_rPr)
        target_r = target_run._element
        if target_r.rPr is not None:
            target_r.remove(target_r.rPr)
        target_r.insert(0, new_rPr)


def set_drawingml_rpr_element_font(a_rPr_element, font_name: str):
    if a_rPr_element is None or not font_name:
        return
    for local_name, tag in (("latin", "a:latin"), ("ea", "a:ea"), ("cs", "a:cs")):
        el = a_rPr_element.find(f"{{{A_NS}}}{local_name}")
        if el is None:
            el = OxmlElement(tag)
            a_rPr_element.append(el)
        el.set("typeface", font_name)


def set_drawingml_r_element_font(a_r_element, font_name: str):
    if a_r_element is None or not font_name:
        return
    a_rPr = a_r_element.find(f"{{{A_NS}}}rPr")
    if a_rPr is None:
        a_rPr = OxmlElement("a:rPr")
        a_r_element.insert(0, a_rPr)
    set_drawingml_rpr_element_font(a_rPr, font_name)


def set_pptx_run_font(run, font_name: str):
    if not run or not font_name:
        return
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r = run._r
        if hasattr(r, "get_or_add_rPr"):
            rPr = r.get_or_add_rPr()
        else:
            rPr = r.find(f"{{{A_NS}}}rPr")
            if rPr is None:
                rPr = OxmlElement("a:rPr")
                r.insert(0, rPr)
        set_drawingml_rpr_element_font(rPr, font_name)
    except Exception:
        pass


def set_pptx_paragraph_default_font(paragraph, font_name: str):
    if not paragraph or not font_name:
        return
    try:
        p = paragraph._p
        pPr = p.find(f"{{{A_NS}}}pPr")
        if pPr is None:
            pPr = OxmlElement("a:pPr")
            p.insert(0, pPr)
        defRPr = pPr.find(f"{{{A_NS}}}defRPr")
        if defRPr is None:
            defRPr = OxmlElement("a:defRPr")
            pPr.append(defRPr)
        set_drawingml_rpr_element_font(defRPr, font_name)
    except Exception:
        pass


def _apply_ppt_paragraph_translation(paragraph, normalized: str, translated: str):
    result_norm = append_translation_to_original(normalized, translated) if APPEND_TRANSLATION else translated
    result_text = _ppt_denormalize_linebreaks(result_norm)
    first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    original_font_name = first_run.font.name
    original_size = first_run.font.size
    original_bold = first_run.font.bold
    original_italic = first_run.font.italic
    original_underline = first_run.font.underline

    first_run.text = result_text
    for r in paragraph.runs[1:]:
        r.text = ""

    target_font_name = get_target_font_name()
    if target_font_name:
        set_pptx_run_font(first_run, target_font_name)
        set_pptx_paragraph_default_font(paragraph, target_font_name)
    elif original_font_name:
        first_run.font.name = original_font_name
        first_run.font.size = original_size
        first_run.font.bold = original_bold
        first_run.font.italic = original_italic
        first_run.font.underline = original_underline


def _translate_ppt_paragraphs_batch(paragraphs, allow_english: bool, protect: bool):
    if not paragraphs:
        return
    jobs = []
    for paragraph in paragraphs:
        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            continue
        normalized = _ppt_normalize_linebreaks(full_text)
        if allow_english:
            if not (_text_contains_korean(normalized) or _text_contains_english(normalized)):
                continue
        else:
            if not _text_contains_korean(normalized):
                continue

        lines = normalized.split("\n")
        line_pos = []
        line_texts = []
        for pos, line in enumerate(lines):
            if not line.strip():
                continue
            if _text_contains_korean(line) or (allow_english and _text_contains_english(line)):
                line_pos.append(pos)
                line_texts.append(line)
        if not line_texts:
            continue

        jobs.append((paragraph, normalized, lines, line_pos))
    if not jobs:
        return

    if protect:
        protected_list = []
        mappings = []
        for paragraph, normalized, lines, line_pos in jobs:
            for pos in line_pos:
                protected, mapping = _protect_tokens(lines[pos])
                protected_list.append(protected)
                mappings.append(mapping)
        translated_protected = get_translation_batch(protected_list)
        translated_list = [_restore_tokens(t, m) for t, m in zip(translated_protected, mappings)]
    else:
        to_translate = []
        for paragraph, normalized, lines, line_pos in jobs:
            for pos in line_pos:
                to_translate.append(lines[pos])
        translated_list = get_translation_batch(to_translate)

    cursor = 0
    for paragraph, normalized, lines, line_pos in jobs:
        count = len(line_pos)
        segs = translated_list[cursor : cursor + count]
        cursor += count
        out_lines = list(lines)
        for pos, seg in zip(line_pos, segs):
            out_lines[pos] = seg
        translated = "\n".join(out_lines)
        original_texts.append(normalized)
        translated_texts.append(translated)
        _apply_ppt_paragraph_translation(paragraph, normalized, translated)


def translate_shape_for_ppt(shape):
    if shape.has_text_frame:
        _translate_ppt_paragraphs_batch(list(shape.text_frame.paragraphs), allow_english=True, protect=True)
    elif shape.shape_type == 6:
        for sub_shape in shape.shapes:
            translate_shape_for_ppt(sub_shape)
    elif shape.has_table:
        paragraphs = []
        for row in shape.table.rows:
            for cell in row.cells:
                paragraphs.extend(list(cell.text_frame.paragraphs))
        _translate_ppt_paragraphs_batch(paragraphs, allow_english=False, protect=False)


def translate_ppt(input_file: str, output_file: str):
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        print(f"[进度] 正在翻译 PPT: 第 {i}/{total_slides} 页...")
        for shape in slide.shapes:
            translate_shape_for_ppt(shape)
    prs.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] PPT 已保存: {output_file}")
    return output_file


def translate_excel_xlsx(input_file: str, output_file: str):
    wb = load_workbook(input_file, keep_vba=True)
    existing_sheet_names = set()
    original_sheet_names = list(wb.sheetnames)
    translated_sheet_names = get_translation_batch(original_sheet_names) if original_sheet_names else []
    total_sheets = len(original_sheet_names)

    for i, sheet_name in enumerate(original_sheet_names, 1):
        print(f"[进度] 正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})...")
        translated_sheet_name = translated_sheet_names[i - 1] if i - 1 < len(translated_sheet_names) else get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            unique_sheet_name = f"{cleaned_sheet_name[:27]}_{count}"[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)

        ws = wb[sheet_name]
        ws.title = unique_sheet_name

        row_jobs = []
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            jobs = []
            texts = []
            for cell in row:
                if cell.value is None:
                    continue
                if getattr(cell, "data_type", None) == "f":
                    continue
                if not isinstance(cell.value, str):
                    continue
                if _looks_like_excel_formula(cell.value):
                    continue
                normalized_value = _normalize_text_newlines(str(cell.value))
                if not normalized_value.strip():
                    continue
                if not _text_contains_korean(normalized_value):
                    continue

                lines = normalized_value.split("\n")
                line_pos = []
                line_texts = []
                for pos, line in enumerate(lines):
                    if line.strip() and _text_contains_korean(line):
                        line_pos.append(pos)
                        line_texts.append(line)
                if not line_texts:
                    continue

                offset = len(texts)
                jobs.append(
                    {
                        "cell": cell,
                        "orig": normalized_value,
                        "lines": lines,
                        "line_pos": line_pos,
                        "offset": offset,
                        "count": len(line_texts),
                        "font": copy.copy(cell.font),
                        "border": copy.copy(cell.border),
                        "align": copy.copy(cell.alignment),
                        "fill": copy.copy(cell.fill),
                    }
                )
                texts.extend(line_texts)
            if jobs:
                row_jobs.append((row_idx, jobs, texts))

        for row_idx, jobs, texts in row_jobs:
            translated_list = get_translation_batch(texts)
            for job in jobs:
                segs = translated_list[job["offset"] : job["offset"] + job["count"]]
                out_lines = list(job["lines"])
                for pos, seg in zip(job["line_pos"], segs):
                    out_lines[pos] = seg
                merged_trans = "\n".join(out_lines)
                original_texts.append(job["orig"])
                translated_texts.append(merged_trans)
                job["cell"].value = append_translation_to_original(job["orig"], merged_trans, job["cell"]) if APPEND_TRANSLATION else merged_trans
                if TARGET_FONT_NAME:
                    job["cell"].font = Font(
                        name=TARGET_FONT_NAME,
                        size=job["font"].size,
                        bold=job["font"].bold,
                        italic=job["font"].italic,
                        underline=job["font"].underline,
                        color=job["font"].color,
                    )
                else:
                    job["cell"].font = job["font"]
                job["cell"].border = job["border"]
                job["cell"].alignment = Alignment(
                    horizontal=job["align"].horizontal,
                    vertical="center",
                    text_rotation=job["align"].text_rotation,
                    wrap_text=True,
                    shrink_to_fit=job["align"].shrink_to_fit,
                    indent=job["align"].indent,
                )
                job["cell"].fill = job["fill"]

    wb.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] Excel(.xlsx) 已保存: {output_file}")
    return output_file


def translate_excel_xls(input_file: str, output_file: str):
    base_output = os.path.splitext(output_file)[0]
    output_file_xlsx = base_output + ".xlsx"
    excel_file = pd.ExcelFile(input_file)
    writer = pd.ExcelWriter(output_file_xlsx, engine="openpyxl")
    existing_sheet_names = set()
    original_sheet_names = list(excel_file.sheet_names)
    translated_sheet_names = get_translation_batch(original_sheet_names) if original_sheet_names else []
    total_sheets = len(original_sheet_names)

    for i, sheet_name in enumerate(original_sheet_names, 1):
        print(f"[进度] 正在翻译 Excel(.xls): 第 {i}/{total_sheets} 个工作表 ({sheet_name})...")
        translated_sheet_name = translated_sheet_names[i - 1] if i - 1 < len(translated_sheet_names) else get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            unique_sheet_name = f"{cleaned_sheet_name[:27]}_{count}"[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)

        df = pd.read_excel(input_file, sheet_name=sheet_name)
        row_jobs = []
        for idx in df.index:
            row_texts = []
            cell_jobs = []
            for col in df.columns:
                cell_value = df.at[idx, col]
                if not (pd.notna(cell_value) and isinstance(cell_value, str)):
                    continue
                if _looks_like_excel_formula(cell_value):
                    continue
                normalized_value = _normalize_text_newlines(str(cell_value))
                if not normalized_value.strip():
                    continue
                if not _text_contains_korean(normalized_value):
                    continue

                lines = normalized_value.split("\n")
                line_pos = []
                line_texts = []
                for pos, line in enumerate(lines):
                    if line.strip() and _text_contains_korean(line):
                        line_pos.append(pos)
                        line_texts.append(line)
                if not line_texts:
                    continue

                offset = len(row_texts)
                row_texts.extend(line_texts)
                cell_jobs.append(
                    {
                        "col": col,
                        "orig": normalized_value,
                        "lines": lines,
                        "line_pos": line_pos,
                        "offset": offset,
                        "count": len(line_texts),
                    }
                )
            if cell_jobs:
                row_jobs.append((idx, row_texts, cell_jobs))

        for idx, row_texts, cell_jobs in row_jobs:
            translated_list = get_translation_batch(row_texts)
            for job in cell_jobs:
                segs = translated_list[job["offset"] : job["offset"] + job["count"]]
                out_lines = list(job["lines"])
                for pos, seg in zip(job["line_pos"], segs):
                    out_lines[pos] = seg
                merged_trans = "\n".join(out_lines)
                original_texts.append(job["orig"])
                translated_texts.append(merged_trans)
                df.at[idx, job["col"]] = append_translation_to_original(job["orig"], merged_trans) if APPEND_TRANSLATION else merged_trans

        df.to_excel(writer, sheet_name=unique_sheet_name, index=False)

    writer.close()
    if TARGET_FONT_NAME:
        try:
            wb = load_workbook(output_file_xlsx)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            cell.font = Font(
                                name=TARGET_FONT_NAME,
                                size=cell.font.size,
                                bold=cell.font.bold,
                                italic=cell.font.italic,
                                underline=cell.font.underline,
                                color=cell.font.color,
                            )
            wb.save(output_file_xlsx)
        except Exception:
            pass
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] Excel(.xls -> .xlsx) 已保存: {output_file_xlsx}")
    return output_file_xlsx


def translate_word(input_file: str, output_file: str):
    doc = Document(input_file)
    target_font_name = get_target_font_name()
    processed_elements = set()

    def set_style_font(style_element, font_name: str):
        if style_element is None or not font_name:
            return
        rPr = style_element.find(qn("w:rPr"))
        if rPr is None:
            rPr = OxmlElement("w:rPr")
            style_element.append(rPr)
        rFonts = rPr.find(qn("w:rFonts"))
        if rFonts is None:
            rFonts = OxmlElement("w:rFonts")
            rPr.append(rFonts)
        rFonts.set(qn("w:ascii"), font_name)
        rFonts.set(qn("w:hAnsi"), font_name)
        rFonts.set(qn("w:eastAsia"), font_name)
        rFonts.set(qn("w:cs"), font_name)
        for k in (qn("w:asciiTheme"), qn("w:hAnsiTheme"), qn("w:eastAsiaTheme"), qn("w:csTheme")):
            try:
                rFonts.attrib.pop(k, None)
            except Exception:
                pass

    pending_texts = []

    def collect_texts_from_container(container_el):
        for p_el in container_el.xpath(".//w:p"):
            try:
                full_text = Paragraph(p_el, doc).text
                if full_text and full_text.strip():
                    for line in full_text.replace("\v", "\n").split("\n"):
                        if line.strip():
                            pending_texts.append(line.strip())
            except Exception:
                continue
        for p_el in container_el.xpath(".//a:p"):
            try:
                t_nodes = p_el.xpath(".//a:t")
                combined = "".join([n.text for n in t_nodes if n.text])
                if combined.strip():
                    pending_texts.append(combined.strip())
            except Exception:
                continue

    collect_texts_from_container(doc.element.body)
    for section in doc.sections:
        for h in [section.header, section.first_page_header, section.even_page_header]:
            if h:
                collect_texts_from_container(h._element)
        for f in [section.footer, section.first_page_footer, section.even_page_footer]:
            if f:
                collect_texts_from_container(f._element)

    if pending_texts:
        unique_texts = list(dict.fromkeys(pending_texts))
        to_translate = [
            t
            for t in unique_texts
            if _text_contains_korean(t) and _translation_cache.get((DEFAULT_SRC_LANG, DEFAULT_TGT_LANG, t)) is None
        ]
        if to_translate:
            print(f"[系统] Word 预扫描发现 {len(to_translate)} 条新文本，开始批量翻译...")
            get_translation_batch(to_translate)

    def translate_word_paragraph(paragraph):
        p_id = id(paragraph._element)
        if p_id in processed_elements:
            return
        processed_elements.add(p_id)

        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            return

        try:
            in_table = bool(paragraph._element.xpath("ancestor::w:tbl"))
        except Exception:
            in_table = False

        template_run = paragraph.runs[0] if paragraph.runs else None
        lines = [l.strip() for l in full_text.replace("\v", "\n").splitlines() if l.strip()]
        if not lines:
            return

        if APPEND_TRANSLATION:
            translated_lines = []
            for l in lines:
                if in_table:
                    if _text_contains_korean(l):
                        t_line = _translate_text_with_protection(l)
                        translated_lines.append(t_line)
                        original_texts.append(l)
                        translated_texts.append(t_line)
                    else:
                        translated_lines.append(l)
                else:
                    if _text_contains_korean(l) or _text_contains_english(l):
                        t_line = _translate_text_with_protection(l)
                        translated_lines.append(t_line)
                        original_texts.append(l)
                        translated_texts.append(t_line)
                    else:
                        translated_lines.append(l)
            p_el = paragraph._element
            for r in p_el.xpath("./w:r"):
                p_el.remove(r)
            for i, (orig, trans) in enumerate(zip(lines, translated_lines)):
                run_orig = paragraph.add_run(orig)
                if template_run:
                    apply_run_format(template_run, run_orig)
                run_orig.add_break(WD_BREAK.LINE)
                run_trans = paragraph.add_run(trans)
                if template_run:
                    apply_run_format(template_run, run_trans)
                if target_font_name:
                    set_docx_run_font(run_trans, target_font_name)
                if i < len(lines) - 1:
                    run_trans.add_break(WD_BREAK.LINE)
        else:
            normalized = "\n".join(lines)
            if in_table:
                if not _text_contains_korean(normalized):
                    return
                t = _translate_preserve_newlines_with_protection(normalized, allow_english=False)
            else:
                if not (_text_contains_korean(normalized) or _text_contains_english(normalized)):
                    return
                t = _translate_preserve_newlines_with_protection(normalized, allow_english=True)
            original_texts.append(normalized)
            translated_texts.append(t)
            if paragraph.runs:
                first_run = paragraph.runs[0]
                first_run.text = t
                for r in paragraph.runs[1:]:
                    try:
                        has_drawing = bool(r.element.xpath(".//w:drawing"))
                    except Exception:
                        has_drawing = False
                    if not has_drawing:
                        r.text = ""
                if target_font_name:
                    set_docx_run_font(first_run, target_font_name)
            else:
                new_run = paragraph.add_run(t)
                if target_font_name:
                    set_docx_run_font(new_run, target_font_name)

    def fill_translations_in_container(container_el):
        for p_el in container_el.xpath(".//w:p"):
            try:
                translate_word_paragraph(Paragraph(p_el, doc))
            except Exception:
                continue
        for p_el in container_el.xpath(".//a:p"):
            p_id = id(p_el)
            if p_id in processed_elements:
                continue
            processed_elements.add(p_id)
            t_nodes = p_el.xpath(".//a:t")
            if not t_nodes:
                continue
            combined = "".join([n.text for n in t_nodes if n.text])
            if not combined.strip():
                continue
            if not (_text_contains_korean(combined) or _text_contains_english(combined)):
                continue
            t = _translate_text_with_protection(combined)
            original_texts.append(combined)
            translated_texts.append(t)
            if APPEND_TRANSLATION:
                t_nodes[0].text = combined
                for n in t_nodes[1:]:
                    n.text = ""
                r_nodes = p_el.xpath(".//a:r")
                if r_nodes:
                    last_r = r_nodes[0]
                    br = OxmlElement("a:br")
                    last_r.addnext(br)
                    new_r = OxmlElement("a:r")
                    new_t = OxmlElement("a:t")
                    new_t.text = t
                    new_r.append(new_t)
                    br.addnext(new_r)
                    if target_font_name:
                        set_drawingml_r_element_font(new_r, target_font_name)
            else:
                t_nodes[0].text = t
                for n in t_nodes[1:]:
                    n.text = ""

    fill_translations_in_container(doc.element.body)
    for section in doc.sections:
        for h in [section.header, section.first_page_header, section.even_page_header]:
            if h:
                fill_translations_in_container(h._element)
        for f in [section.footer, section.first_page_footer, section.even_page_footer]:
            if f:
                fill_translations_in_container(f._element)

    if target_font_name:
        try:
            for s in doc.styles:
                try:
                    if getattr(s, "font", None) is not None:
                        s.font.name = target_font_name
                except Exception:
                    pass
                try:
                    set_style_font(getattr(s, "_element", None), target_font_name)
                except Exception:
                    pass
            if not APPEND_TRANSLATION:
                for r in doc.element.xpath(".//w:r"):
                    set_docx_r_element_font(r, target_font_name)
            for r in doc.element.xpath(".//a:r"):
                set_drawingml_r_element_font(r, target_font_name)
        except Exception:
            pass

    doc.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] Word 已保存: {output_file}")
    return output_file


def build_default_output_path(input_file: str) -> str:
    base, ext = os.path.splitext(input_file)
    if ext.lower() == ".xls":
        return f"{base}_translated.xlsx"
    return f"{base}_translated{ext}"


def translate_file(input_file: str, output_file: str | None = None) -> str:
    if not output_file:
        output_file = build_default_output_path(input_file)
    ext = os.path.splitext(input_file)[1].lower()
    reset_runtime_buffers()
    load_nllb_model()
    if ext in {".ppt", ".pptx"}:
        return translate_ppt(input_file, output_file)
    if ext == ".xlsx":
        return translate_excel_xlsx(input_file, output_file)
    if ext == ".xls":
        return translate_excel_xls(input_file, output_file)
    if ext == ".docx":
        return translate_word(input_file, output_file)
    raise ValueError(f"暂不支持的文件类型: {ext}")


def quick_translate(file_path: str | None = None):
    if not file_path:
        file_path = DEFAULT_INPUT_FILE
    return translate_file(file_path)


if __name__ == "__main__":
    print("=" * 60)
    print("NLLB 3.3B 文件翻译工具 - Colab 版本 (Excel / PPT / Word)")
    print("=" * 60)
    print(f"对照模式(APPEND_TRANSLATION): {bool(APPEND_TRANSLATION)}")
    print(f"同步生成语料库(GENERATE_CORPUS): {bool(GENERATE_CORPUS)}")
    print(f"Corpus 目录: {CORPUS_DIR}")
    print(f"输入文件: {DEFAULT_INPUT_FILE}")
    if os.path.exists(DEFAULT_INPUT_FILE):
        output_file = quick_translate(DEFAULT_INPUT_FILE)
        print(f"[完成] 输出文件: {output_file}")
    else:
        print(f"[错误] 文件不存在: {DEFAULT_INPUT_FILE}")
        print("请修改 DEFAULT_INPUT_FILE，或调用: quick_translate('/path/to/file')")
